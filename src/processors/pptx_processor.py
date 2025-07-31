"""
PPTX document processor - Improved version
"""

import copy
from typing import Any, Dict, List

from pptx import Presentation

from processors.base import BaseDocumentProcessor


class PPTXProcessor(BaseDocumentProcessor):
    """Processor for PPTX documents"""

    def extract_text(self, file_path: str) -> Dict[str, Any]:
        """Extract text content from PPTX while preserving structure and shape metadata"""
        prs = Presentation(file_path)
        content = {"slides": []}

        for slide_num, slide in enumerate(prs.slides):
            slide_content = {"slide_number": slide_num + 1, "shapes": []}

            for shape_idx, shape in enumerate(slide.shapes):
                shape_data = {
                    "shape_index": shape_idx,
                    "shape_type": str(shape.shape_type),
                    "shape_id": shape.shape_id if hasattr(shape, "shape_id") else None,
                    "text_content": None,
                    "has_text_frame": False,
                    "text_frames": [],
                }

                # Handle shapes with text
                if hasattr(shape, "text_frame") and shape.text_frame:
                    shape_data["has_text_frame"] = True
                    for para_idx, paragraph in enumerate(shape.text_frame.paragraphs):
                        para_text = paragraph.text.strip()
                        if para_text:
                            shape_data["text_frames"].append(
                                {
                                    "paragraph_index": para_idx,
                                    "text": para_text,
                                    "original_text": para_text,  # Keep original for reference
                                }
                            )

                # Handle shapes with direct text property
                elif hasattr(shape, "text") and shape.text.strip():
                    shape_data["text_content"] = shape.text.strip()

                # Handle table shapes
                elif hasattr(shape, "table"):
                    table_data = []
                    for row_idx, row in enumerate(shape.table.rows):
                        row_data = []
                        for cell_idx, cell in enumerate(row.cells):
                            cell_text = cell.text.strip()
                            if cell_text:
                                row_data.append(
                                    {
                                        "cell_index": cell_idx,
                                        "text": cell_text,
                                        "original_text": cell_text,
                                    }
                                )
                            else:
                                row_data.append(
                                    {
                                        "cell_index": cell_idx,
                                        "text": "",
                                        "original_text": "",
                                    }
                                )
                        table_data.append({"row_index": row_idx, "cells": row_data})
                    shape_data["table_data"] = table_data

                # Only add shapes that have text content
                if (
                    shape_data["text_content"]
                    or shape_data["text_frames"]
                    or shape_data.get("table_data")
                ):
                    slide_content["shapes"].append(shape_data)

            content["slides"].append(slide_content)

        return content

    def reconstruct_document(
        self, original_path: str, translated_content: Dict[str, Any], output_path: str
    ) -> str:
        """Reconstruct PPTX document by modifying original file and replacing text with translations"""
        # Load original presentation to preserve all formatting
        prs = Presentation(original_path)

        # Process each slide
        for slide_data in translated_content["slides"]:
            slide_num = slide_data["slide_number"] - 1  # Convert to 0-based index

            if slide_num >= len(prs.slides):
                continue  # Skip if slide doesn't exist in original

            slide = prs.slides[slide_num]

            # Process each shape with translated content
            for shape_data in slide_data["shapes"]:
                shape_idx = shape_data["shape_index"]

                if shape_idx >= len(slide.shapes):
                    continue  # Skip if shape doesn't exist

                shape = slide.shapes[shape_idx]

                # Handle text frame content
                if shape_data["has_text_frame"] and shape_data["text_frames"]:
                    if hasattr(shape, "text_frame") and shape.text_frame:
                        # Clear existing paragraphs if we have translated content
                        # shape.text_frame.clear()

                        # Add translated paragraphs
                        for idx, frame_data in enumerate(shape_data["text_frames"]):
                            print(frame_data)
                            if frame_data["text"].strip():
                                p = shape.text_frame.paragraphs[idx]
                                p.text = frame_data["text"]
                                # Note: This approach maintains basic formatting but may lose some advanced styling

                # Handle direct text content
                elif shape_data["text_content"]:
                    if hasattr(shape, "text"):
                        shape.text = shape_data["text_content"]

                # Handle table content
                elif shape_data.get("table_data"):
                    if hasattr(shape, "table"):
                        table = shape.table
                        for row_data in shape_data["table_data"]:
                            row_idx = row_data["row_index"]
                            if row_idx < len(table.rows):
                                row = table.rows[row_idx]
                                for cell_data in row_data["cells"]:
                                    cell_idx = cell_data["cell_index"]
                                    if cell_idx < len(row.cells) and cell_data["text"]:
                                        row.cells[cell_idx].text = cell_data["text"]

        # Save the modified presentation
        prs.save(output_path)
        return output_path

    def get_translatable_texts(self, extracted_content: Dict[str, Any]) -> List[str]:
        """
        Extract all translatable texts from the extracted content
        Param:
            extracted_content (Dict[str, Any]): The content structure extracted from the PPTX file
        Returns:
            List[str]: List of all translatable texts found in the content"""
        texts = []

        for slide_data in extracted_content["slides"]:
            for shape_data in slide_data["shapes"]:
                # Get text from text frames
                if shape_data["text_frames"]:
                    for frame_data in shape_data["text_frames"]:
                        if frame_data["text"].strip():
                            texts.append(frame_data["text"])

                # Get direct text content
                elif shape_data["text_content"]:
                    texts.append(shape_data["text_content"])
                # Get table text
                elif shape_data.get("table_data"):
                    for row_data in shape_data["table_data"]:
                        for cell_data in row_data["cells"]:
                            if cell_data["text"].strip():
                                texts.append(cell_data["text"])
        return texts
    
    def get_pairs_translation(self, translated_content: Dict[str, Any]) -> Dict[str, str]:
        """
        Extract original->translated pairs from the translated content structure
        Param:
            translated_content (Dict[str, Any]): The translated content structure
        Returns:
            Dict[str, str]: Dictionary of original text to translated text pairs
        """
        pairs = {}
        for slide_data in translated_content["slides"]:
            for shape_data in slide_data["shapes"]:
                # Handle text frames
                if shape_data["text_frames"]:
                    for frame_data in shape_data["text_frames"]:
                        if frame_data["original_text"] and frame_data["text"].strip():
                            pairs[frame_data["original_text"]] = frame_data["text"]

                # Handle direct text content
                elif shape_data["text_content"]:
                    if shape_data.get("original_text") and shape_data["text_content"].strip():
                        pairs[shape_data["original_text"]] = shape_data["text_content"]

                # Handle table data
                elif shape_data.get("table_data"):
                    for row_data in shape_data["table_data"]:
                        for cell_data in row_data["cells"]:
                            if cell_data.get("original_text") and cell_data["text"].strip():
                                pairs[cell_data["original_text"]] = cell_data["text"]
        
        return pairs

    def apply_translations(
        self, extracted_content: Dict[str, Any], translations: List[str]
    ) -> Dict[str, Any]:
        """
        Apply translations to the extracted content structure
        Param:
            extracted_content (Dict[str, Any]): The original extracted content structure
            translations (List[str]): List of translated texts to apply
        Returns:
            Dict[str, Any]: The translated content structure with applied translations
        """
        translated_content = copy.deepcopy(extracted_content)
        translation_idx = 0

        for slide_data in translated_content["slides"]:
            for shape_data in slide_data["shapes"]:
                # Apply translations to text frames
                if shape_data["text_frames"]:
                    for frame_data in shape_data["text_frames"]:
                        if frame_data["text"].strip() and translation_idx < len(
                            translations
                        ):
                            frame_data["text"] = translations[translation_idx]
                            translation_idx += 1

                # Apply translation to direct text content
                elif shape_data["text_content"] and translation_idx < len(translations):
                    shape_data["text_content"] = translations[translation_idx]
                    translation_idx += 1

                # Apply translations to table content
                elif shape_data.get("table_data"):
                    for row_data in shape_data["table_data"]:
                        for cell_data in row_data["cells"]:
                            if cell_data["text"].strip() and translation_idx < len(
                                translations
                            ):
                                cell_data["text"] = translations[translation_idx]
                                translation_idx += 1
        return translated_content
